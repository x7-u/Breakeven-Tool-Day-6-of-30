"""Day 6. PowerPoint export for the BREAK CVP tool.

Three slides:
  1. Cover. Company, period, BE units, BE revenue, CM ratio.
  2. Break-even chart (PNG full-bleed).
  3. Sensitivity tornado + AI verdict + actions.

python-pptx is optional. is_available() returns False when it is not
installed; the server skips the PPTX path gracefully in that case.
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False

from cvp_chart import render_breakeven_png, render_tornado_png
from cvp_maths import CVPResult


def is_available() -> bool:
    return _HAS_PPTX


INK = (0x0F, 0x11, 0x17)
PURPLE = (0x7C, 0x3A, 0xED)
GREY = (0x6B, 0x72, 0x80)


def _ccy_sym(c: str) -> str:
    return {"GBP": "£", "USD": "$", "EUR": "€"}.get(c.upper(), "")


def write_pptx(result: CVPResult, out_path: Path) -> Path:
    if not _HAS_PPTX:
        raise RuntimeError("python-pptx is not installed.")
    md = result.cvp.metadata
    inp = result.cvp.inputs
    h = result.cvp.headline
    sym = _ccy_sym(md.currency)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: cover
    s1 = prs.slides.add_slide(blank)
    _add_text(s1, "BREAK", Inches(0.6), Inches(0.5), Inches(4), Inches(0.6),
              size=22, bold=True, color=PURPLE)
    _add_text(s1, "Day 06 . Cost-Volume-Profit", Inches(0.6), Inches(1.0),
              Inches(8), Inches(0.4), size=12, color=GREY)
    _add_text(s1, md.company, Inches(0.6), Inches(1.7), Inches(12), Inches(1.0),
              size=36, bold=True, color=INK)
    _add_text(s1, f"{md.period_label}  |  {md.currency}", Inches(0.6),
              Inches(2.7), Inches(12), Inches(0.5), size=14, color=GREY)
    if h.break_even_units == float("inf"):
        be_line = "Break-even: NEVER (contribution margin <= 0)"
    else:
        be_line = (f"Break-even: {h.break_even_units:,.0f} units  |  "
                   f"{sym}{h.break_even_revenue:,.0f} revenue  |  "
                   f"CM ratio {h.cm_ratio:.1%}")
    _add_text(s1, be_line, Inches(0.6), Inches(3.6), Inches(12), Inches(0.6),
              size=18, color=INK)
    if result.commentary and result.commentary.headline:
        _add_text(s1, result.commentary.headline, Inches(0.6), Inches(4.5),
                  Inches(12), Inches(2), size=14, color=INK)

    # Slide 2: break-even chart
    s2 = prs.slides.add_slide(blank)
    _add_text(s2, "Break-even chart", Inches(0.6), Inches(0.4), Inches(12),
              Inches(0.5), size=20, bold=True, color=INK)
    png = render_breakeven_png(
        inp, title=md.company, currency=md.currency,
        current_volume=md.current_volume,
    )
    s2.shapes.add_picture(BytesIO(png), Inches(0.6), Inches(1.2),
                          width=Inches(12), height=Inches(5.6))

    # Slide 3: tornado + AI commentary
    s3 = prs.slides.add_slide(blank)
    _add_text(s3, "Sensitivity & verdict", Inches(0.6), Inches(0.4),
              Inches(12), Inches(0.5), size=20, bold=True, color=INK)
    if result.cvp.tornado:
        png = render_tornado_png(result.cvp.tornado, title="BE swing per variable")
        s3.shapes.add_picture(BytesIO(png), Inches(0.6), Inches(1.2),
                              width=Inches(7), height=Inches(5.0))
    if result.commentary and not result.commentary.skipped:
        c = result.commentary
        x = Inches(8.0)
        _add_text(s3, c.headline or "", x, Inches(1.2), Inches(4.8),
                  Inches(0.8), size=14, bold=True, color=INK)
        _add_text(s3, c.summary or "", x, Inches(2.1), Inches(4.8),
                  Inches(2.5), size=11, color=INK)
        if c.actions:
            actions_text = "Actions:\n" + "\n".join(f"- {a}" for a in c.actions)
            _add_text(s3, actions_text, x, Inches(4.8), Inches(4.8),
                      Inches(2.0), size=11, color=GREY)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def _add_text(slide, text, left, top, width, height, *,
              size=14, bold=False, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.color.rgb = RGBColor(*color)
    return box
